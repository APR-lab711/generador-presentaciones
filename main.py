from fastapi import FastAPI
import openai
from pptx import Presentation

# Reemplaza con tu clave de OpenAI
openai.api_key = "TU_API_KEY"

app = FastAPI()

@app.get("/generar_presentacion/")
def generar_presentacion(tema: str):
    # Generar contenido con la IA
    respuesta = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role": "system", "content": "Eres un experto en presentaciones."},
                  {"role": "user", "content": f"Crea una presentación sobre {tema}"}]
    )
    contenido = respuesta["choices"][0]["message"]["content"].split("\n")

    # Crear un archivo PowerPoint con el contenido
    prs = Presentation()
    for slide_text in contenido[:5]:  # Limitar a 5 diapositivas
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title, content = slide.shapes.title, slide.placeholders[1]
        title.text, content.text = slide_text.split(":")[0], slide_text.split(":")[1]

    prs.save("presentacion.pptx")
    return {"mensaje": "Presentación generada", "archivo": "presentacion.pptx"}
